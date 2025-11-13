#!/usr/bin/env python3
"""
BN Overseas DevOps Presentation Report Generator
Generates a comprehensive PowerPoint presentation with project and DevOps information
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime

class PresentationGenerator:
    def __init__(self, output_file="BN_Overseas_DevOps_Presentation.pptx"):
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        self.output_file = output_file
        
        # Color scheme
        self.primary_color = RGBColor(25, 118, 210)  # Blue
        self.secondary_color = RGBColor(56, 142, 60)  # Green
        self.accent_color = RGBColor(251, 140, 0)    # Orange
        self.text_dark = RGBColor(33, 33, 33)
        self.text_light = RGBColor(255, 255, 255)
    
    def add_title_slide(self):
        """Add title slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.primary_color
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_p = title_frame.paragraphs[0]
        title_p.text = "BN Overseas"
        title_p.font.size = Pt(66)
        title_p.font.bold = True
        title_p.font.color.rgb = self.text_light
        title_p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        sub_p = subtitle_frame.paragraphs[0]
        sub_p.text = "Study Abroad Platform - DevOps Architecture & Deployment"
        sub_p.font.size = Pt(28)
        sub_p.font.color.rgb = self.text_light
        sub_p.alignment = PP_ALIGN.CENTER
        
        # Date
        date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
        date_frame = date_box.text_frame
        date_p = date_frame.paragraphs[0]
        date_p.text = f"November 2025"
        date_p.font.size = Pt(16)
        date_p.font.color.rgb = self.text_light
        date_p.alignment = PP_ALIGN.CENTER
    
    def add_content_slide(self, title, content_points):
        """Add a content slide with bullet points"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(44)
        title_p.font.bold = True
        title_p.font.color.rgb = self.primary_color
        
        # Add line under title
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(9), Inches(0))
        line.line.color.rgb = self.accent_color
        line.line.width = Pt(3)
        
        # Content
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.4), Inches(5.5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        for i, point in enumerate(content_points):
            if i > 0:
                text_frame.add_paragraph()
            p = text_frame.paragraphs[i]
            p.text = point
            p.font.size = Pt(20)
            p.font.color.rgb = self.text_dark
            p.level = 0
            p.space_before = Pt(8)
            p.space_after = Pt(8)
    
    def add_two_column_slide(self, title, left_title, left_points, right_title, right_points):
        """Add a two-column comparison slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(44)
        title_p.font.bold = True
        title_p.font.color.rgb = self.primary_color
        
        # Add line under title
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(9), Inches(0))
        line.line.color.rgb = self.accent_color
        line.line.width = Pt(3)
        
        # Left column
        left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.2), Inches(0.5))
        left_title_frame = left_title_box.text_frame
        left_title_p = left_title_frame.paragraphs[0]
        left_title_p.text = left_title
        left_title_p.font.size = Pt(24)
        left_title_p.font.bold = True
        left_title_p.font.color.rgb = self.secondary_color
        
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(4.2), Inches(5))
        left_frame = left_box.text_frame
        left_frame.word_wrap = True
        
        for i, point in enumerate(left_points):
            if i > 0:
                left_frame.add_paragraph()
            p = left_frame.paragraphs[i]
            p.text = f"• {point}"
            p.font.size = Pt(16)
            p.font.color.rgb = self.text_dark
            p.space_before = Pt(4)
            p.space_after = Pt(4)
        
        # Right column
        right_title_box = slide.shapes.add_textbox(Inches(5.3), Inches(1.5), Inches(4.2), Inches(0.5))
        right_title_frame = right_title_box.text_frame
        right_title_p = right_title_frame.paragraphs[0]
        right_title_p.text = right_title
        right_title_p.font.size = Pt(24)
        right_title_p.font.bold = True
        right_title_p.font.color.rgb = self.accent_color
        
        right_box = slide.shapes.add_textbox(Inches(5.3), Inches(2.1), Inches(4.2), Inches(5))
        right_frame = right_box.text_frame
        right_frame.word_wrap = True
        
        for i, point in enumerate(right_points):
            if i > 0:
                right_frame.add_paragraph()
            p = right_frame.paragraphs[i]
            p.text = f"• {point}"
            p.font.size = Pt(16)
            p.font.color.rgb = self.text_dark
            p.space_before = Pt(4)
            p.space_after = Pt(4)
    
    def generate(self):
        """Generate the complete presentation"""
        
        # Slide 1: Title
        self.add_title_slide()
        
        # Slide 2: Project Overview
        self.add_content_slide(
            "Project Overview",
            [
                "🎓 BN Overseas - Study Abroad Platform",
                "Built with Next.js 14, React, and TypeScript",
                "Serves students seeking international education opportunities",
                "Comprehensive course management and appointment booking system",
                "Real-time test preparation and progress tracking",
                "Multi-role authentication (Student, Instructor, Admin)"
            ]
        )
        
        # Slide 3: Key Features
        self.add_two_column_slide(
            "Key Features",
            "Student Features",
            [
                "Course enrollment & browsing",
                "Appointment booking with advisors",
                "Practice tests & performance tracking",
                "Country information guides",
                "Blog & educational resources",
                "Personal dashboard"
            ],
            "Admin Features",
            [
                "User & role management",
                "Content management (Blog, Services)",
                "Appointment management",
                "Course administration",
                "Payment tracking",
                "Analytics & statistics"
            ]
        )
        
        # Slide 4: Technology Stack
        self.add_two_column_slide(
            "Technology Stack",
            "Frontend",
            [
                "Next.js 14 (App Router)",
                "React & TypeScript",
                "Tailwind CSS",
                "Shadcn/UI Components",
                "React Hook Form",
                "Zod Validation"
            ],
            "Backend & Database",
            [
                "Next.js API Routes",
                "NextAuth.js (Authentication)",
                "Prisma ORM",
                "PostgreSQL",
                "AWS S3 (File Storage)",
                "Razorpay (Payments)"
            ]
        )
        
        # Slide 5: DevOps Architecture
        self.add_content_slide(
            "DevOps Architecture",
            [
                "🏗️ Infrastructure as Code (Terraform)",
                "   • AWS VPC, EC2, Security Groups",
                "   • Automated provisioning and scaling",
                "",
                "⚙️ Configuration Management (Ansible)",
                "   • Automated deployment pipeline",
                "   • Environment variable management",
                "",
                "📊 Monitoring (Nagios)",
                "   • Real-time server monitoring",
                "   • CPU, memory, disk, HTTP status alerts"
            ]
        )
        
        # Slide 6: AWS Infrastructure
        self.add_content_slide(
            "AWS Infrastructure Setup",
            [
                "☁️ Region: Mumbai (ap-south-1)",
                "",
                "🖥️ EC2 Instances:",
                "   • App Server: t3.medium (Next.js + Nginx + PM2)",
                "   • Nagios Server: t3.small (Monitoring)",
                "",
                "🔐 Security:",
                "   • VPC with public subnet",
                "   • Security groups for SSH, HTTP, NRPE",
                "   • SSH keypair authentication"
            ]
        )
        
        # Slide 7: Terraform
        self.add_content_slide(
            "Terraform - Infrastructure as Code",
            [
                "✓ Automated VPC & Network Setup",
                "   • VPC (10.0.0.0/16), Public Subnet, IGW",
                "",
                "✓ Security Groups",
                "   • App: SSH (restricted), HTTP/HTTPS (open)",
                "   • Nagios: SSH (restricted), NRPE (monitored)",
                "",
                "✓ EC2 Instance Configuration",
                "   • Ubuntu 22.04 AMI",
                "   • Key pair injection",
                "   • Public IP assignment",
                "",
                "✓ Outputs: IPs for Ansible inventory"
            ]
        )
        
        # Slide 8: Ansible - Deployment
        self.add_content_slide(
            "Ansible - Automated Deployment",
            [
                "📦 Application Deployment:",
                "   • Node.js 18.x installation",
                "   • GitHub repo cloning",
                "   • NPM dependency installation",
                "   • Next.js build compilation",
                "",
                "🚀 Service Management:",
                "   • PM2 process manager setup",
                "   • Nginx reverse proxy configuration",
                "   • Auto-restart on reboot",
                "",
                "🔒 Environment Configuration:",
                "   • .env.local creation from variables",
                "   • Database, auth, and API keys",
                "   • AWS and payment gateway credentials"
            ]
        )
        
        # Slide 9: Nagios - Monitoring
        self.add_content_slide(
            "Nagios - Monitoring & Alerting",
            [
                "📈 Metrics Monitored:",
                "   • HTTP service availability (port 80)",
                "   • CPU load and usage",
                "   • Disk space utilization",
                "   • Memory consumption",
                "",
                "🔔 Alert Mechanisms:",
                "   • Real-time notifications",
                "   • Web UI dashboard",
                "   • Status visualization",
                "",
                "🔗 Integration:",
                "   • NRPE (Nagios Remote Plugin Executor)",
                "   • Agent-based monitoring on app server"
            ]
        )
        
        # Slide 10: Deployment Pipeline
        self.add_content_slide(
            "Deployment Pipeline (Step-by-Step)",
            [
                "1️⃣ AWS Credentials & SSH Setup",
                "   aws configure, ssh-keygen",
                "",
                "2️⃣ Terraform Apply",
                "   terraform init → terraform apply",
                "",
                "3️⃣ Ansible Configuration",
                "   Update inventory with IPs, create vars.yml",
                "",
                "4️⃣ Deploy Application",
                "   ansible-playbook playbook.yml -e @vars.yml",
                "",
                "5️⃣ Setup Monitoring",
                "   ansible-playbook nagios-playbook.yml"
            ]
        )
        
        # Slide 11: File Structure
        self.add_content_slide(
            "Project Structure",
            [
                "📁 terraform/",
                "   • main.tf: VPC, EC2, Security Groups",
                "   • variables.tf: Configuration variables",
                "",
                "📁 ansible/",
                "   • playbook.yml: App deployment",
                "   • nagios-playbook.yml: Monitoring setup",
                "   • inventory.ini: Host inventory",
                "   • vars.yml: Environment variables",
                "",
                "📁 nagios/",
                "   • README.md: Monitoring reference",
                "",
                "📄 DEPLOYMENT.md: Complete guide with commands"
            ]
        )
        
        # Slide 12: Environment Variables
        self.add_content_slide(
            "Environment Configuration",
            [
                "🔐 Database:",
                "   DATABASE_URL, connection pooling",
                "",
                "🔑 Authentication:",
                "   NEXTAUTH_SECRET, JWT secrets",
                "",
                "📧 Email Service:",
                "   SMTP configuration for notifications",
                "",
                "💳 Payment Gateways:",
                "   Razorpay and Stripe credentials",
                "",
                "☁️ Cloud Services:",
                "   AWS S3, Zoom API, Twilio"
            ]
        )
        
        # Slide 13: Security Considerations
        self.add_content_slide(
            "Security Best Practices",
            [
                "🔒 SSH Access Control",
                "   • Restrict SSH to specific IP in terraform",
                "   • Use strong key-based authentication",
                "",
                "🛡️ Environment Secrets",
                "   • Never commit .env to GitHub",
                "   • Store credentials in vars.yml (gitignored)",
                "",
                "🔐 AWS IAM",
                "   • Use least-privilege IAM policies",
                "   • Rotate access keys regularly",
                "",
                "🌐 HTTPS",
                "   • Setup Let's Encrypt SSL certificates",
                "   • Configure Nginx for SSL"
            ]
        )
        
        # Slide 14: Monitoring & Alerts
        self.add_two_column_slide(
            "Monitoring & Alerting Strategy",
            "Metrics Tracked",
            [
                "HTTP response time",
                "Server availability",
                "CPU utilization",
                "Disk space usage",
                "Memory consumption",
                "Network traffic"
            ],
            "Alert Actions",
            [
                "Email notifications",
                "Dashboard updates",
                "Service restart triggers",
                "Log aggregation",
                "Escalation policies",
                "Historical trending"
            ]
        )
        
        # Slide 15: Scaling & Production
        self.add_content_slide(
            "Production Enhancements",
            [
                "🚀 Scaling Options:",
                "   • Load Balancer (ALB) for multiple instances",
                "   • Auto Scaling Groups for elastic capacity",
                "",
                "🐳 Containerization:",
                "   • Docker for consistent environments",
                "   • ECS/Fargate for serverless deployment",
                "",
                "📊 Advanced Monitoring:",
                "   • Prometheus + Grafana for visualization",
                "   • ELK Stack for log aggregation",
                "",
                "🔄 CI/CD Pipeline:",
                "   • GitHub Actions for automated testing",
                "   • Automated deployments on git push"
            ]
        )
        
        # Slide 16: Cost Optimization
        self.add_content_slide(
            "Cost Optimization",
            [
                "💰 Instance Sizing:",
                "   • t3.medium for app (burstable, cost-effective)",
                "   • t3.small for Nagios monitoring",
                "",
                "📈 Auto-scaling:",
                "   • Scale up during peak hours",
                "   • Scale down during off-peak periods",
                "",
                "🗄️ Storage Optimization:",
                "   • S3 lifecycle policies for old logs",
                "   • RDS backup optimization",
                "",
                "🌍 Regional Selection:",
                "   • Mumbai (ap-south-1) for India operations",
                "   • Lower latency, cost-effective"
            ]
        )
        
        # Slide 17: Troubleshooting Guide
        self.add_content_slide(
            "Troubleshooting Common Issues",
            [
                "❌ Ansible Connection Failed",
                "   → Check security group SSH rules, key permissions",
                "",
                "❌ Next.js Build Error",
                "   → Verify all dependencies, check logs: pm2 logs nextjs-app",
                "",
                "❌ Database Connection Error",
                "   → Validate DATABASE_URL, network connectivity",
                "",
                "❌ Nagios Not Monitoring",
                "   → Update app IP in /usr/local/nagios/etc/servers/app.cfg",
                "",
                "❌ Nginx Not Proxying",
                "   → Check config syntax: nginx -t, view logs"
            ]
        )
        
        # Slide 18: Maintenance & Operations
        self.add_content_slide(
            "Maintenance & Operations",
            [
                "🔄 Regular Backups:",
                "   • Database backups (daily/weekly)",
                "   • Configuration backups",
                "",
                "📝 Logging:",
                "   • Application logs: pm2 logs",
                "   • Nginx logs: /var/log/nginx/",
                "   • System logs: /var/log/syslog",
                "",
                "🔐 Security Updates:",
                "   • OS patches and kernel updates",
                "   • Node.js and npm dependency updates",
                "",
                "📊 Performance Monitoring:",
                "   • Monitor Nagios dashboard daily",
                "   • Review application metrics"
            ]
        )
        
        # Slide 19: Quick Commands Reference
        self.add_content_slide(
            "Quick Commands Reference",
            [
                "# Terraform",
                "terraform init  |  terraform plan  |  terraform apply  |  terraform destroy",
                "",
                "# Ansible",
                "ansible-playbook -i inventory.ini playbook.yml -e @vars.yml",
                "",
                "# SSH Access",
                "ssh -i ~/.ssh/deploy-key ubuntu@APP_PUBLIC_IP",
                "",
                "# PM2 Management",
                "pm2 status  |  pm2 logs nextjs-app  |  pm2 restart nextjs-app",
                "",
                "# Nagios Access",
                "http://NAGIOS_PUBLIC_IP (username: nagios, password: nagios123)"
            ]
        )
        
        # Slide 20: Conclusion
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.secondary_color
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_p = title_frame.paragraphs[0]
        title_p.text = "Complete DevOps Infrastructure"
        title_p.font.size = Pt(54)
        title_p.font.bold = True
        title_p.font.color.rgb = self.text_light
        title_p.alignment = PP_ALIGN.CENTER
        
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(2))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        sub_p = subtitle_frame.paragraphs[0]
        sub_p.text = "Terraform • Ansible • Nagios • AWS"
        sub_p.font.size = Pt(32)
        sub_p.font.color.rgb = self.text_light
        sub_p.alignment = PP_ALIGN.CENTER
        
        # Save presentation
        self.prs.save(self.output_file)
        print(f"✅ Presentation generated: {self.output_file}")


if __name__ == "__main__":
    generator = PresentationGenerator()
    generator.generate()
